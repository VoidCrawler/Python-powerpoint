from pptx import Presentation
import os
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import PP_ALIGN
import pandas as pd

df_xlsx = pd.read_excel("D:/Upthrust/Frank/PoC Cookbook Automation/Excel/Volkswagen Cookbook Excel.xlsx")

x = 0
# the row
y = 0
# the column



z = df_xlsx.iloc[x,y]

prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]


title.text = "VW Cookbook"
subtitle.text = "Growthmarketing"


def main():


    z = df_xlsx.iloc[x, y]
    normal_slide = prs.slide_layouts[5]
    slide_2 = prs.slides.add_slide(normal_slide)
    title = slide_2.shapes.title


    title.text = "Experiment 1: Search Campaign Per Model"



    for shape in slide_2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame


    text_frame = shape.text_frame
    text_frame.clear()



    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(25.3)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme



    txBox = slide_2.shapes.add_textbox(Inches(0), Inches(6.5), Inches(2), Inches(1))
    tf = txBox.text_frame



    for shape in slide_2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame


    text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Started: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme

    y = y + 1
    z = str(df_xlsx.iloc[x,y])


    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme

    p = tf.add_paragraph()



    run = p.add_run()
    run.text = 'Status: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme

    y = y + 1
    z = df_xlsx.iloc[x,y]

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme


    p = tf.add_paragraph()

    run = p.add_run()
    run.text = 'Channel(s): '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    y = y + 1
    z = df_xlsx.iloc[x,y]

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme


    txBox = slide_2.shapes.add_textbox(Inches(3), Inches(6.5), Inches(2), Inches(1))
    tf = txBox.text_frame

    for shape in slide_2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame


    text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Current Results: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme

    txBox = slide_2.shapes.add_textbox(Inches(5), Inches(6.5), Inches(2), Inches(1))
    tf = txBox.text_frame


    for shape in slide_2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame


    text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Total Reach: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    text_frame = shape.text_frame

    y = y + 1
    z = str(df_xlsx.iloc[x,y])

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme

    p = tf.add_paragraph()

    text_frame = shape.text_frame


    run = p.add_run()
    run.text = 'Total Clicks: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    text_frame = shape.text_frame

    y = y + 1
    z = str(df_xlsx.iloc[x,y])


    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme

    p = tf.add_paragraph()

    text_frame = shape.text_frame


    run = p.add_run()
    run.text = 'Media Spend: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    text_frame = shape.text_frame

    y = y + 1
    z = str(df_xlsx.iloc[x,y])

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme


    txBox = slide_2.shapes.add_textbox(Inches(7.5), Inches(6.5), Inches(2), Inches(1))
    tf = txBox.text_frame


    for shape in slide_2.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        # do things with the text frame


    text_frame = shape.text_frame

    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Total Leads: '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    text_frame = shape.text_frame

    y = y + 1
    z = str(df_xlsx.iloc[x,y])

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme

    p = tf.add_paragraph()

    text_frame = shape.text_frame


    run = p.add_run()
    run.text = 'Cost Per Lead: : '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme


    text_frame = shape.text_frame

    y = y + 1
    z = str(df_xlsx.iloc[x,y])

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme


    p = tf.add_paragraph()

    text_frame = shape.text_frame


    run = p.add_run()
    run.text = 'Ads Set Up: : '

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme

    y = y + 1
    z = str(df_xlsx.iloc[x,y])

    text_frame = shape.text_frame

    run = p.add_run()
    run.text = z

    font = run.font
    font.name = 'Century Goth'
    font.size = Pt(14)
    font.bold = False
    font.italic = None  # cause value to be inherited from theme
    prs.save('D:/Upthrust/Frank/PoC Cookbook Automation/Powerpoints/Test.pptx')

    x = x + 1
    y = 0
    z = str(df_xlsx.iloc[x,y])
    main()

os.startfile("D:/Upthrust/Frank/PoC Cookbook Automation/Powerpoints/Test.pptx")


y = y + 1

z = df_xlsx.iloc[x,y]


